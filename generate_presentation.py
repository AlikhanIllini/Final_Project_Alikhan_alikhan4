#!/usr/bin/env python3
"""
Generate professional PowerPoint presentation for Stock Cards
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Define colors
BLUE = RGBColor(59, 130, 246)
GREEN = RGBColor(16, 185, 129)
ORANGE = RGBColor(245, 158, 11)
RED = RGBColor(239, 68, 68)
DARK_GRAY = RGBColor(31, 41, 55)
LIGHT_GRAY = RGBColor(249, 250, 251)

def add_title_slide(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Main title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "STOCK CARDS"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = BLUE
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(14), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Your Personal Stock Organization System"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = DARK_GRAY
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(14), Inches(0.7))
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = "📈 From Spreadsheet Chaos to Visual Clarity"
    tagline_para = tagline_frame.paragraphs[0]
    tagline_para.font.size = Pt(28)
    tagline_para.font.color.rgb = GREEN
    tagline_para.alignment = PP_ALIGN.CENTER

    # Author info
    author_box = slide.shapes.add_textbox(Inches(1), Inches(7.5), Inches(14), Inches(1))
    author_frame = author_box.text_frame
    author_text = "Alikhan\nINFO Semester Project • December 2025\n\ngithub.com/AlikhanIllini/Final_Project_Alikhan_alikhan4"
    author_frame.text = author_text
    author_para = author_frame.paragraphs[0]
    author_para.font.size = Pt(16)
    author_para.font.color.rgb = DARK_GRAY
    author_para.alignment = PP_ALIGN.CENTER

def add_problem_slide(prs):
    """Slide 2: The Problem"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "THE PROBLEM"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = BLUE
    title_para.alignment = PP_ALIGN.CENTER

    # Line
    line = slide.shapes.add_shape(1, Inches(3), Inches(1.3), Inches(10), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()

    # Content
    content_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(12), Inches(4))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    problems = [
        "📊  Cluttered spreadsheets",
        "🤯  Information overload",
        "❌  No personalization",
        "⚠️  Hard to prioritize what matters"
    ]

    for i, problem in enumerate(problems):
        p = content_frame.add_paragraph() if i > 0 else content_frame.paragraphs[0]
        p.text = problem
        p.font.size = Pt(36)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(20)
        p.alignment = PP_ALIGN.LEFT

    # Who needs this
    who_box = slide.shapes.add_textbox(Inches(2), Inches(6.5), Inches(12), Inches(1.5))
    who_frame = who_box.text_frame
    who_frame.text = "Who needs better?\nStudents • Young investors • Professionals"
    who_para = who_frame.paragraphs[0]
    who_para.font.size = Pt(28)
    who_para.font.italic = True
    who_para.font.color.rgb = ORANGE
    who_para.alignment = PP_ALIGN.CENTER

def add_solution_slide(prs):
    """Slide 3: The Solution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "STOCK CARDS"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = BLUE
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(14), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "A visual, card-based stock organizer"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = DARK_GRAY
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Features list
    content_box = slide.shapes.add_textbox(Inches(2.5), Inches(2.5), Inches(11), Inches(4))
    content_frame = content_box.text_frame

    features = [
        "Each stock = One card with:",
        "• Live prices 📊",
        "• Your notes 📝",
        "• Priority levels 🎯",
        "• Custom tags 🏷️",
        "• Price history 📈"
    ]

    for i, feature in enumerate(features):
        p = content_frame.add_paragraph() if i > 0 else content_frame.paragraphs[0]
        p.text = feature
        if i == 0:
            p.font.size = Pt(32)
            p.font.bold = True
        else:
            p.font.size = Pt(30)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(15)
        p.alignment = PP_ALIGN.LEFT

    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(2), Inches(6.8), Inches(12), Inches(1))
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = '"Think Trello for your stock watchlist"'
    tagline_para = tagline_frame.paragraphs[0]
    tagline_para.font.size = Pt(32)
    tagline_para.font.italic = True
    tagline_para.font.color.rgb = GREEN
    tagline_para.alignment = PP_ALIGN.CENTER

def add_features_slide(prs):
    """Slide 4: Key Features"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "WHAT IT DOES"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = BLUE
    title_para.alignment = PP_ALIGN.CENTER

    # Features in 2x3 grid
    features = [
        ("🔐", "Multi-User System", "Secure • Private"),
        ("📊", "Live Stock Prices", "Auto-fetch • Cached"),
        ("🏷️", "Smart Tagging", "Color-coded • Custom"),
        ("🔍", "Powerful Filtering", "Search • Sort • Filter"),
        ("📈", "Price Intelligence", "7-day/30-day changes"),
        ("📧", "Weekly Digest", "Email summaries")
    ]

    x_positions = [Inches(1.5), Inches(8.5)]
    y_positions = [Inches(2), Inches(4), Inches(6)]

    for i, (emoji, title, desc) in enumerate(features):
        row = i % 3
        col = i // 3
        x = x_positions[col]
        y = y_positions[row]

        # Feature box
        box = slide.shapes.add_textbox(x, y, Inches(6), Inches(1.5))
        frame = box.text_frame
        frame.word_wrap = True

        # Emoji
        p1 = frame.paragraphs[0]
        p1.text = emoji
        p1.font.size = Pt(40)
        p1.space_after = Pt(5)

        # Title
        p2 = frame.add_paragraph()
        p2.text = title
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = DARK_GRAY
        p2.space_after = Pt(3)

        # Description
        p3 = frame.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(18)
        p3.font.color.rgb = RGBColor(107, 114, 128)

def add_tech_slide(prs):
    """Slide 5: Tech & Vision"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "BUILT WITH DJANGO"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = BLUE
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(14), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Deployed Live • 5 Data Models • RESTful API"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(22)
    subtitle_para.font.color.rgb = DARK_GRAY
    subtitle_para.alignment = PP_ALIGN.CENTER

    # What I Built (Left column)
    built_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(6), Inches(4))
    built_frame = built_box.text_frame

    p = built_frame.paragraphs[0]
    p.text = "What I Built:"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = GREEN

    built_items = [
        "✓ User authentication",
        "✓ Real-time price integration",
        "✓ Database architecture",
        "✓ Responsive interface",
        "✓ Production deployment"
    ]

    for item in built_items:
        p = built_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(22)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(10)

    # Future Vision (Right column)
    future_box = slide.shapes.add_textbox(Inches(8.5), Inches(2.5), Inches(6), Inches(4))
    future_frame = future_box.text_frame

    p = future_frame.paragraphs[0]
    p.text = "Future Vision:"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ORANGE

    future_items = [
        "→ Interactive charts",
        "→ Portfolio analytics",
        "→ Real-time updates",
        "→ Mobile app"
    ]

    for item in future_items:
        p = future_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(22)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(10)

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(7.5), Inches(14), Inches(1))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Alikhan | github.com/AlikhanIllini/Final_Project_Alikhan_alikhan4"
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(16)
    footer_para.font.color.rgb = RGBColor(107, 114, 128)
    footer_para.alignment = PP_ALIGN.CENTER

def add_demo_slide(prs):
    """Slide 6: Demo slide (for backup)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "LIVE DEMONSTRATION"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = BLUE
    title_para.alignment = PP_ALIGN.CENTER

    # Demo steps
    content_box = slide.shapes.add_textbox(Inches(3), Inches(2.5), Inches(10), Inches(5))
    content_frame = content_box.text_frame

    steps = [
        "1. Dashboard with stock cards",
        "2. Create new stock (NVDA)",
        "3. View card details & price history",
        "4. Tag management",
        "5. Filters & sorting"
    ]

    for i, step in enumerate(steps):
        p = content_frame.add_paragraph() if i > 0 else content_frame.paragraphs[0]
        p.text = step
        p.font.size = Pt(32)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(20)
        p.alignment = PP_ALIGN.LEFT

    # Note
    note_box = slide.shapes.add_textbox(Inches(2), Inches(7.5), Inches(12), Inches(0.8))
    note_frame = note_box.text_frame
    note_frame.text = "(Use this slide only if live demo fails - otherwise skip to live demo)"
    note_para = note_frame.paragraphs[0]
    note_para.font.size = Pt(18)
    note_para.font.italic = True
    note_para.font.color.rgb = RED
    note_para.alignment = PP_ALIGN.CENTER

# Generate all slides
print("Creating presentation...")
add_title_slide(prs)
add_problem_slide(prs)
add_solution_slide(prs)
add_features_slide(prs)
add_tech_slide(prs)
add_demo_slide(prs)

# Save presentation
output_file = "Stock_Cards_Presentation.pptx"
prs.save(output_file)
print(f"✓ Presentation saved as: {output_file}")
print(f"\nNow convert to PDF:")
print(f"1. Open {output_file} in PowerPoint/Keynote")
print(f"2. File → Export → PDF")
print(f"3. Or use Google Slides: Upload → Download as PDF")

